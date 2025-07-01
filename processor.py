from fastapi import FastAPI, File, UploadFile
import cv2
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

def extract_frames(video_path):
    frames = []
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise ValueError("Unable to open video file")

    frame_count = 0
    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        
        frame_count += 1
        if frame_count % 30 == 0:
            frames.append(frame)

    cap.release()
    return frames

def add_text_to_frame(frame, text, position=(50, 50)):
    font = cv2.FONT_HERSHEY_SIMPLEX
    font_scale = 1
    font_color = (255, 255, 255)
    line_type = 2

    cv2.putText(frame, text, position, font, font_scale, font_color, line_type)
    return frame

def create_presentation(frames, output_folder):
    prs = Presentation()
    for i, frame in enumerate(frames):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Slide {i+1}"

        img_path = os.path.join(output_folder, f"frame_{i}.jpg")
        cv2.imwrite(img_path, frame)
        
        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(img_path, left, top, height=Inches(5.5))
        
        os.remove(img_path)  # Удаляем временный файл

    presentation_path = os.path.join(output_folder, "presentation.pptx")
    prs.save(presentation_path)
    return presentation_path

app = FastAPI()

@app.post("/process/")
async def process_video(file: UploadFile = File(...)):
    file_location = f"temp_{file.filename}"
    with open(file_location, "wb+") as file_object:
        file_object.write(file.file.read())

    frames = extract_frames(file_location)
    
    output_folder = f"output_{file.filename}"
    os.makedirs(output_folder, exist_ok=True)

    for i, frame in enumerate(frames):
        text = f"Frame {i+1}"
        frame_with_text = add_text_to_frame(frame, text)
        frames[i] = frame_with_text

    presentation_path = create_presentation(frames, output_folder)

    os.remove(file_location)

    return {"presentation_path": presentation_path}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8001)
